"""
RGUKT Basar Web Crawler and Document Ingestion Agent
Crawls https://www.rgukt.ac.in/ and https://hub.rgukt.ac.in/
Extracts content, downloads documents, and indexes into Chroma vector store
"""

import os
import re
import json
import time
import hashlib
import requests
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse
from datetime import datetime
from typing import Dict, List, Tuple, Optional
import logging

# Document extraction libraries
import pdfplumber
from pypdf import PdfReader
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# ============================================================
# CONFIGURATION
# ============================================================
BASE_DOMAINS = [
    "https://www.rgukt.ac.in",
    "https://hub.rgukt.ac.in"
]

ALLOWED_EXTENSIONS = {
    '.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.zip'
}

CHUNK_SIZE = 1000  # characters
CHUNK_OVERLAP = 150  # characters

# Metadata storage
METADATA_FILE = "crawl_metadata.json"
DOCUMENTS_DIR = "rgukt_documents"
CHUNKS_FILE = "rgukt_chunks.jsonl"

# ============================================================
# UTILITY FUNCTIONS
# ============================================================
def is_valid_url(url: str) -> bool:
    """Check if URL belongs to allowed domains"""
    try:
        parsed = urlparse(url)
        return any(parsed.netloc == urlparse(domain).netloc for domain in BASE_DOMAINS)
    except:
        return False

def is_document_url(url: str) -> bool:
    """Check if URL points to a document"""
    parsed = urlparse(url)
    path = parsed.path.lower()
    return any(path.endswith(ext) for ext in ALLOWED_EXTENSIONS)

def get_file_extension(url: str) -> str:
    """Get file extension from URL"""
    parsed = urlparse(url)
    path = parsed.path.lower()
    for ext in ALLOWED_EXTENSIONS:
        if path.endswith(ext):
            return ext
    return ""

def generate_id(content: str) -> str:
    """Generate unique ID from content"""
    return hashlib.md5(content.encode()).hexdigest()[:16]

def sanitize_filename(filename: str) -> str:
    """Sanitize filename for safe storage"""
    return re.sub(r'[<>:"/\\|?*]', '_', filename)

# ============================================================
# DOCUMENT EXTRACTION
# ============================================================
def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file"""
    text = ""
    try:
        # Try pdfplumber first (better for complex PDFs)
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n\n"
    except Exception as e1:
        try:
            # Fallback to pypdf
            reader = PdfReader(file_path)
            for page in reader.pages:
                text += page.extract_text() + "\n\n"
        except Exception as e2:
            logger.error(f"Error extracting PDF {file_path}: {e1}, {e2}")
    return text

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file"""
    text = ""
    try:
        doc = DocxDocument(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        logger.error(f"Error extracting DOCX {file_path}: {e}")
    return text

def extract_text_from_doc(file_path: str) -> str:
    """Extract text from DOC file (limited support)"""
    # DOC files are binary and hard to parse without antiword
    # Return empty for now, can be enhanced later
    return ""

def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from XLSX file"""
    text = ""
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text += f"\n--- Sheet: {sheet} ---\n"
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text += row_text + "\n"
    except Exception as e:
        logger.error(f"Error extracting XLSX {file_path}: {e}")
    return text

def extract_text_from_xls(file_path: str) -> str:
    """Extract text from XLS file"""
    # XLS files require xlrd, can be added if needed
    return ""

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file"""
    text = ""
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {i} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        logger.error(f"Error extracting PPTX {file_path}: {e}")
    return text

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPT file"""
    # PPT files require different library, can be added if needed
    return ""

def extract_text_from_zip(file_path: str) -> str:
    """Extract text from ZIP file (list contents)"""
    import zipfile
    text = ""
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            text += "ZIP Contents:\n"
            for info in zf.infolist():
                text += f"- {info.filename} ({info.file_size} bytes)\n"
    except Exception as e:
        logger.error(f"Error extracting ZIP {file_path}: {e}")
    return text

def extract_text_from_document(file_path: str, url: str) -> str:
    """Extract text from document based on file type"""
    ext = get_file_extension(url)
    
    extractors = {
        '.pdf': extract_text_from_pdf,
        '.docx': extract_text_from_docx,
        '.doc': extract_text_from_doc,
        '.xlsx': extract_text_from_xlsx,
        '.xls': extract_text_from_xls,
        '.pptx': extract_text_from_pptx,
        '.ppt': extract_text_from_pptx,
        '.zip': extract_text_from_zip,
    }
    
    extractor = extractors.get(ext)
    if extractor:
        return extractor(file_path)
    return ""

# ============================================================
# WEB CRAWLER
# ============================================================
class RGUKTCrawler:
    def __init__(self, base_urls: List[str], delay: float = 1.0):
        self.base_urls = base_urls
        self.delay = delay  # Delay between requests (seconds)
        self.visited_urls = set()
        self.urls_to_visit = set(base_urls)
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        })
        
        # Create documents directory
        os.makedirs(DOCUMENTS_DIR, exist_ok=True)
        
        # Load existing metadata
        self.metadata = self.load_metadata()
    
    def load_metadata(self) -> Dict:
        """Load existing crawl metadata"""
        if os.path.exists(METADATA_FILE):
            try:
                with open(METADATA_FILE, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except:
                pass
        return {"documents": {}, "urls": {}}
    
    def save_metadata(self):
        """Save crawl metadata"""
        with open(METADATA_FILE, 'w', encoding='utf-8') as f:
            json.dump(self.metadata, f, indent=2, ensure_ascii=False)
    
    def get_url_hash(self, url: str) -> str:
        """Generate hash for URL"""
        return hashlib.md5(url.encode()).hexdigest()
    
    def is_duplicate(self, url: str, content_hash: str) -> bool:
        """Check if URL/content is duplicate"""
        url_hash = self.get_url_hash(url)
        if url_hash in self.metadata["urls"]:
            existing_hash = self.metadata["urls"][url_hash]
            return existing_hash == content_hash
        return False
    
    def crawl(self, max_pages: int = 1000):
        """Main crawl loop"""
        pages_crawled = 0
        
        while self.urls_to_visit and pages_crawled < max_pages:
            url = self.urls_to_visit.pop()
            
            if url in self.visited_urls:
                continue
            
            logger.info(f"Crawling: {url}")
            
            try:
                response = self.session.get(url, timeout=15)
                response.raise_for_status()
                
                content_type = response.headers.get('Content-Type', '')
                
                # Handle documents
                if is_document_url(url):
                    self.process_document(url, response.content)
                # Handle HTML pages
                elif 'text/html' in content_type:
                    self.process_html_page(url, response.text)
                    # Extract new links
                    new_links = self.extract_links(url, response.text)
                    self.urls_to_visit.update(new_links)
                
                self.visited_urls.add(url)
                pages_crawled += 1
                
                # Be polite
                time.sleep(self.delay)
                
            except Exception as e:
                logger.error(f"Error crawling {url}: {e}")
                self.visited_urls.add(url)
        
        logger.info(f"Crawl complete. Total pages: {pages_crawled}")
        self.save_metadata()
    
    def extract_links(self, base_url: str, html: str) -> set:
        """Extract all internal links from HTML"""
        links = set()
        try:
            soup = BeautifulSoup(html, 'html.parser')
            for a_tag in soup.find_all('a', href=True):
                href = a_tag['href']
                # Convert relative URLs to absolute
                full_url = urljoin(base_url, href)
                
                # Only include internal links
                if is_valid_url(full_url) and full_url not in self.visited_urls:
                    # Skip anchors and javascript
                    if not full_url.startswith('#') and not full_url.startswith('javascript:'):
                        links.add(full_url)
        except Exception as e:
            logger.error(f"Error extracting links from {base_url}: {e}")
        
        return links
    
    def process_html_page(self, url: str, html: str):
        """Process HTML page and extract content"""
        try:
            soup = BeautifulSoup(html, 'html.parser')
            
            # Remove script, style, nav, footer
            for tag in soup.find_all(['script', 'style', 'nav', 'footer', 'header']):
                tag.decompose()
            
            # Extract title
            title = ""
            if soup.title:
                title = soup.title.get_text(strip=True)
            if not title:
                h1 = soup.find('h1')
                if h1:
                    title = h1.get_text(strip=True)
            
            # Extract main content
            body = soup.find('body') or soup
            content = body.get_text(separator='\n', strip=True)
            
            # Clean up content
            lines = [line.strip() for line in content.split('\n') if line.strip()]
            content = '\n'.join(lines)
            
            # Generate content hash
            content_hash = self.get_url_hash(content)
            
            # Check for duplicates
            if self.is_duplicate(url, content_hash):
                logger.info(f"Skipping duplicate: {url}")
                return
            
            # Extract metadata
            metadata = self.extract_metadata(url, title, content, "html")
            
            # Save to metadata store
            url_hash = self.get_url_hash(url)
            self.metadata["documents"][url_hash] = metadata
            self.metadata["urls"][url_hash] = content_hash
            
            # Save raw HTML for reference
            self.save_raw_content(url, html, content, "html")
            
        except Exception as e:
            logger.error(f"Error processing HTML page {url}: {e}")
    
    def process_document(self, url: str, content: bytes):
        """Process and download document"""
        try:
            # Generate filename
            parsed = urlparse(url)
            filename = sanitize_filename(os.path.basename(parsed.path))
            if not filename:
                filename = f"document_{int(time.time())}"
            
            file_path = os.path.join(DOCUMENTS_DIR, filename)
            
            # Save document
            with open(file_path, 'wb') as f:
                f.write(content)
            
            # Extract text
            text = extract_text_from_document(file_path, url)
            
            if not text.strip():
                logger.warning(f"No text extracted from {url}")
                return
            
            # Generate content hash
            content_hash = self.get_url_hash(text)
            
            # Check for duplicates
            if self.is_duplicate(url, content_hash):
                logger.info(f"Skipping duplicate document: {url}")
                os.remove(file_path)  # Remove duplicate file
                return
            
            # Extract metadata
            title = filename
            metadata = self.extract_metadata(url, title, text, "document")
            metadata["file_name"] = filename
            metadata["file_path"] = file_path
            metadata["document_url"] = url
            
            # Save to metadata store
            url_hash = self.get_url_hash(url)
            self.metadata["documents"][url_hash] = metadata
            self.metadata["urls"][url_hash] = content_hash
            
            logger.info(f"Indexed document: {filename}")
            
        except Exception as e:
            logger.error(f"Error processing document {url}: {e}")
    
    def extract_metadata(self, url: str, title: str, content: str, content_type: str) -> Dict:
        """Extract metadata from content"""
        # Try to extract date
        date = self.extract_date(content)
        
        # Try to extract department
        department = self.extract_department(content)
        
        # Determine category
        category = self.categorize_content(url, title, content)
        
        # Extract year
        year = self.extract_year(content, date)
        
        # Generate summary (first 200 chars)
        summary = content[:200].strip() + "..."
        
        # Extract keywords
        keywords = self.extract_keywords(content)
        
        return {
            "id": self.get_url_hash(url),
            "title": title,
            "category": category,
            "department": department,
            "year": year,
            "date": date,
            "source_url": url,
            "document_url": url if content_type == "document" else "",
            "file_name": "",
            "content_type": content_type,
            "summary": summary,
            "keywords": keywords,
            "crawled_at": datetime.now().isoformat()
        }
    
    def extract_date(self, content: str) -> str:
        """Extract date from content"""
        # Common date patterns
        date_patterns = [
            r'\d{1,2}[-/]\d{1,2}[-/]\d{4}',  # DD-MM-YYYY or DD/MM/YYYY
            r'\d{4}[-/]\d{1,2}[-/]\d{1,2}',  # YYYY-MM-DD
            r'\d{1,2}(?:st|nd|rd|th)?\s+(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*[\s,]*\d{4}',
        ]
        
        for pattern in date_patterns:
            match = re.search(pattern, content, re.IGNORECASE)
            if match:
                return match.group(0)
        
        return ""
    
    def extract_year(self, content: str, date: str) -> str:
        """Extract year from content or date"""
        # Try to find 4-digit year
        year_match = re.search(r'\b(20\d{2})\b', content)
        if year_match:
            return year_match.group(1)
        
        # Try from date
        if date:
            year_match = re.search(r'\b(20\d{2})\b', date)
            if year_match:
                return year_match.group(1)
        
        return str(datetime.now().year)
    
    def extract_department(self, content: str) -> str:
        """Extract department from content"""
        departments = [
            "Computer Science", "Electronics", "Mechanical", "Civil",
            "Chemical", "Metallurgical", "Electrical", "Physics",
            "Chemistry", "Mathematics", "Bio Sciences", "Management",
            "Humanities", "Telugu", "English"
        ]
        
        content_lower = content.lower()
        for dept in departments:
            if dept.lower() in content_lower:
                return dept
        
        return "General"
    
    def categorize_content(self, url: str, title: str, content: str) -> str:
        """Categorize content based on URL and content"""
        url_lower = url.lower()
        content_lower = content.lower()
        
        categories = {
            "admission": ["admission", "apply", "counseling", "entrance"],
            "academics": ["academic", "course", "curriculum", "syllabus", "program"],
            "examination": ["exam", "test", "result", "grade", "time table"],
            "departments": ["department", "faculty", "hod"],
            "hostel": ["hostel", "accommodation", "warden"],
            "library": ["library", "book"],
            "placement": ["placement", "job", "career", "recruitment"],
            "scholarship": ["scholarship", "financial", "fee", "concession"],
            "research": ["research", "publication", "patent"],
            "administration": ["administration", "registrar", "dean", "vc"],
            "events": ["event", "workshop", "seminar", "conference"],
            "notices": ["notice", "circular", "announcement"],
            "tenders": ["tender", "bid", "procurement"],
            "contact": ["contact", "address", "phone", "email"],
        }
        
        for category, keywords in categories.items():
            if any(kw in url_lower for kw in keywords):
                return category
            if any(kw in content_lower for kw in keywords):
                return category
        
        return "general"
    
    def extract_keywords(self, content: str, max_keywords: int = 10) -> List[str]:
        """Extract keywords from content"""
        # Simple keyword extraction based on frequency
        words = re.findall(r'\b[a-zA-Z]{4,}\b', content.lower())
        
        # Remove common stop words
        stop_words = {'this', 'that', 'with', 'from', 'have', 'been', 'were', 'was',
                      'they', 'their', 'what', 'when', 'where', 'which', 'about',
                      'more', 'other', 'some', 'than', 'them', 'then', 'these'}
        words = [w for w in words if w not in stop_words]
        
        # Count frequency
        from collections import Counter
        word_freq = Counter(words)
        
        # Return top keywords
        return [word for word, _ in word_freq.most_common(max_keywords)]
    
    def save_raw_content(self, url: str, html: str, text: str, content_type: str):
        """Save raw content for reference"""
        url_hash = self.get_url_hash(url)
        filename = f"{url_hash}.{content_type}"
        filepath = os.path.join(DOCUMENTS_DIR, filename)
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(text)

# ============================================================
# TEXT CHUNKING
# ============================================================
def chunk_text(text: str, metadata: Dict, chunk_size: int = CHUNK_SIZE, 
               overlap: int = CHUNK_OVERLAP) -> List[Dict]:
    """Split text into semantic chunks with metadata"""
    chunks = []
    
    # Split by paragraphs first
    paragraphs = text.split('\n\n')
    
    current_chunk = ""
    chunk_id = 0
    
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        
        # If adding this paragraph exceeds chunk size, save current chunk
        if len(current_chunk) + len(para) > chunk_size and current_chunk:
            # Save chunk
            chunk_metadata = {
                "document_id": metadata["id"],
                "chunk_id": f"{metadata['id']}_chunk_{chunk_id}",
                "title": metadata["title"],
                "category": metadata["category"],
                "year": metadata["year"],
                "source_url": metadata["source_url"],
                "content": current_chunk.strip()
            }
            chunks.append(chunk_metadata)
            chunk_id += 1
            
            # Start new chunk with overlap
            words = current_chunk.split()
            overlap_words = words[-int(overlap/5):]  # Approximate word count
            current_chunk = " ".join(overlap_words) + "\n\n" + para
        else:
            current_chunk += "\n\n" + para if current_chunk else para
    
    # Save final chunk
    if current_chunk.strip():
        chunk_metadata = {
            "document_id": metadata["id"],
            "chunk_id": f"{metadata['id']}_chunk_{chunk_id}",
            "title": metadata["title"],
            "category": metadata["category"],
            "year": metadata["year"],
            "source_url": metadata["source_url"],
            "content": current_chunk.strip()
        }
        chunks.append(chunk_metadata)
    
    return chunks

# ============================================================
# VECTOR STORE INTEGRATION
# ============================================================
def index_chunks_to_vectorstore(chunks: List[Dict], persist_directory: str = "./rgukt2_db"):
    """Index chunks into Chroma vector store"""
    try:
        from langchain_huggingface import HuggingFaceEmbeddings
        from langchain_community.vectorstores import Chroma
        from langchain.schema import Document
        
        logger.info(f"Indexing {len(chunks)} chunks to vector store...")
        
        # Initialize embeddings
        embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
        
        # Convert chunks to LangChain Documents
        documents = []
        for chunk in chunks:
            doc = Document(
                page_content=chunk["content"],
                metadata={
                    "document_id": chunk["document_id"],
                    "chunk_id": chunk["chunk_id"],
                    "title": chunk["title"],
                    "category": chunk["category"],
                    "year": chunk["year"],
                    "source_url": chunk["source_url"]
                }
            )
            documents.append(doc)
        
        # Add to vector store
        vector_store = Chroma(
            persist_directory=persist_directory,
            embedding_function=embeddings
        )
        
        # Add documents in batches
        batch_size = 100
        for i in range(0, len(documents), batch_size):
            batch = documents[i:i+batch_size]
            vector_store.add_documents(batch)
            logger.info(f"Indexed batch {i//batch_size + 1}/{(len(documents)-1)//batch_size + 1}")
        
        logger.info(f"Successfully indexed {len(chunks)} chunks")
        return True
        
    except Exception as e:
        logger.error(f"Error indexing to vector store: {e}")
        return False

# ============================================================
# MAIN CRAWL AND INDEX FUNCTION
# ============================================================
def crawl_and_index(max_pages: int = 500, reindex: bool = False):
    """Main function to crawl and index RGUKT website"""
    logger.info("=" * 60)
    logger.info("Starting RGUKT Basar Web Crawler")
    logger.info("=" * 60)
    
    # Initialize crawler
    crawler = RGUKTCrawler(BASE_DOMAINS, delay=1.0)
    
    # Crawl website
    logger.info(f"Phase 1: Crawling website (max {max_pages} pages)...")
    crawler.crawl(max_pages=max_pages)
    
    # Process and chunk documents
    logger.info("Phase 2: Processing and chunking documents...")
    all_chunks = []
    
    for url_hash, metadata in crawler.metadata["documents"].items():
        # Reconstruct content from saved file
        url = metadata["source_url"]
        content_type = metadata["content_type"]
        
        if content_type == "html":
            filename = f"{url_hash}.html"
        else:
            filename = metadata.get("file_name", f"{url_hash}.txt")
        
        filepath = os.path.join(DOCUMENTS_DIR, filename)
        
        if os.path.exists(filepath):
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            # Chunk the content
            chunks = chunk_text(content, metadata)
            all_chunks.extend(chunks)
            
            logger.info(f"Generated {len(chunks)} chunks from: {metadata['title'][:50]}")
    
    logger.info(f"Total chunks generated: {len(all_chunks)}")
    
    # Save chunks to file
    logger.info("Phase 3: Saving chunks to file...")
    with open(CHUNKS_FILE, 'w', encoding='utf-8') as f:
        for chunk in all_chunks:
            f.write(json.dumps(chunk, ensure_ascii=False) + '\n')
    
    # Index to vector store
    logger.info("Phase 4: Indexing to vector store...")
    success = index_chunks_to_vectorstore(all_chunks)
    
    if success:
        logger.info("=" * 60)
        logger.info("Crawl and indexing complete!")
        logger.info(f"Total documents: {len(crawler.metadata['documents'])}")
        logger.info(f"Total chunks: {len(all_chunks)}")
        logger.info(f"Metadata saved to: {METADATA_FILE}")
        logger.info(f"Chunks saved to: {CHUNKS_FILE}")
        logger.info(f"Documents stored in: {DOCUMENTS_DIR}")
        logger.info("=" * 60)
    else:
        logger.error("Indexing failed. Please check the errors above.")
    
    return crawler.metadata["documents"], all_chunks

# ============================================================
# INCREMENTAL UPDATE
# ============================================================
def incremental_update(max_pages: int = 100):
    """Perform incremental update - only crawl new/modified content"""
    logger.info("Starting incremental update...")
    
    crawler = RGUKTCrawler(BASE_DOMAINS, delay=1.0)
    
    # Only crawl new URLs
    new_urls = crawler.urls_to_visit - crawler.visited_urls
    crawler.urls_to_visit = new_urls
    
    logger.info(f"Found {len(new_urls)} new URLs to crawl")
    
    if new_urls:
        crawler.crawl(max_pages=max_pages)
        
        # Process new documents
        all_chunks = []
        for url_hash, metadata in crawler.metadata["documents"].items():
            if url_hash not in [k for k in crawler.metadata["urls"].keys()]:
                # This is a new document
                url = metadata["source_url"]
                content_type = metadata["content_type"]
                filename = metadata.get("file_name", f"{url_hash}.txt")
                filepath = os.path.join(DOCUMENTS_DIR, filename)
                
                if os.path.exists(filepath):
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                    
                    chunks = chunk_text(content, metadata)
                    all_chunks.extend(chunks)
        
        if all_chunks:
            logger.info(f"Indexing {len(all_chunks)} new chunks...")
            index_chunks_to_vectorstore(all_chunks)
    
    logger.info("Incremental update complete")

# ============================================================
# MAIN ENTRY POINT
# ============================================================
if __name__ == "__main__":
    import sys
    
    if len(sys.argv) > 1:
        if sys.argv[1] == "update":
            # Incremental update
            max_pages = int(sys.argv[2]) if len(sys.argv) > 2 else 100
            incremental_update(max_pages)
        elif sys.argv[1] == "full":
            # Full crawl
            max_pages = int(sys.argv[2]) if len(sys.argv) > 2 else 500
            crawl_and_index(max_pages=max_pages, reindex=True)
        else:
            print("Usage:")
            print("  python crawl_rgukt.py full [max_pages]  - Full crawl and index")
            print("  python crawl_rgukt.py update [max_pages] - Incremental update")
    else:
        # Default: full crawl
        crawl_and_index(max_pages=500)